from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches
import collections

# Define the input file paths
slide1_input_file = "sample_slide1_input.txt"
slide2_input_file = "sample_slide2_input.txt"
font_file = "sample_font_file.ttf"

# Create a new PowerPoint presentation
presentation = Presentation()

# Define a function to set the font style using the provided font file
def set_font_style(run, font_file):
    run.font.name = font_file

# Slide 1
slide1 = presentation.slides.add_slide(presentation.slide_layouts[0])
slide1_title = slide1.shapes.title
slide1_title.text = "Slide 1"

with open(slide1_input_file, "r") as file:
    slide1_content = file.read()

slide1_textbox = slide1.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
slide1_text_frame = slide1_textbox.text_frame
slide1_text_frame.word_wrap = True
slide1_text_frame.text = slide1_content

for paragraph in slide1_text_frame.paragraphs:
    for run in paragraph.runs:
        set_font_style(run, font_file)
        run.font.size = Pt(12)

# Slide 2
slide2 = presentation.slides.add_slide(presentation.slide_layouts[1])
slide2_title = slide2.shapes.title
slide2_title.text = "Slide 2"

with open(slide2_input_file, "r") as file:
    slide2_content = file.read()

slide2_textbox = slide2.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
slide2_text_frame = slide2_textbox.text_frame
slide2_text_frame.word_wrap = True
slide2_text_frame.text = slide2_content

for paragraph in slide2_text_frame.paragraphs:
    for run in paragraph.runs:
        set_font_style(run, font_file)
        run.font.size = Pt(12)

# Save the presentation
output_file = "output.pptx"
presentation.save(output_file)
print(f"Presentation saved as '{output_file}'.")
